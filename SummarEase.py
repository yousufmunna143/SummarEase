import streamlit as st
import requests
import fitz  
from pptx import Presentation 
import docx 
import google.generativeai as genai
# from google.generativeai import GenerateAI
import os

genai.configure(api_key="AIzaSyCP4Q1Iad4GL88kI-HvU-XqRmA5I9dCCpc")

def extract_text_from_pdf(file):
    doc = fitz.open(stream=file.read(), filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_ppt(file):
    prs = Presentation(file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def extract_text_from_docx(file):
    doc = docx.Document(file)
    text = ""
    for para in doc.paragraphs:
        text += para.text
    return text

def extract_text(file, file_type):
    if file_type == "pdf":
        return extract_text_from_pdf(file)
    elif file_type == "ppt":
        return extract_text_from_ppt(file)
    elif file_type == "docx":
        return extract_text_from_docx(file)
    else:
        return "Unsupported file type"
    
generator = genai.GenerativeModel('gemini-1.0-pro')

if "chat" not in st.session_state:
    st.session_state.chat = generator.start_chat(history = [])

if "messages" not in st.session_state:
    st.session_state.messages = [] 

def summarize_text(text):
    prompt = "Please summarize the following text: \n" + text
    response = st.session_state.chat.send_message(prompt) 
    return response.text

st.title("Document Summarizer - SummerEase")

uploaded_file = st.file_uploader("Upload a file", type=["pdf", "ppt", "docx"])

if uploaded_file:
    file_type = uploaded_file.name.split(".")[-1]
    text = extract_text(uploaded_file, file_type)
    
    if text:
        if st.button("Summarize"):
            summary = summarize_text(text)
            st.write("Summary:")
            st.write(summary)
    else:
        st.write("Could not extract text from the uploaded file.")